"""Peek at font sizes in layout placeholders of a .pptx file.

Iterates all slide layouts and prints placeholder metadata including
font size information when available.

Usage:
    python peek_fonts.py                     # auto-find .pptx in uploads/
    python peek_fonts.py path/to/file.pptx   # explicit file
"""

import glob
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt


def emu_to_inches(emu_val):
    """Convert EMU to inches, rounded to 2 decimals."""
    return round(emu_val / 914400, 2)


def emu_to_pt(emu_val):
    """Convert EMU to points (font size unit)."""
    if emu_val is None:
        return None
    return round(emu_val / 12700, 1)


def get_font_info(placeholder):
    """Extract font size from a placeholder using multiple strategies.

    Returns a dict with font size info and which method succeeded.
    python-pptx font.size is in EMU (English Metric Units).
    """
    results = {
        "paragraph_font_size_pt": None,
        "run_font_size_pt": None,
        "defRPr_size_pt": None,
        "source": "none",
        "has_text_frame": False,
        "paragraph_count": 0,
        "sample_text": "",
    }

    if not placeholder.has_text_frame:
        return results

    results["has_text_frame"] = True
    tf = placeholder.text_frame
    results["paragraph_count"] = len(tf.paragraphs)

    for para in tf.paragraphs:
        # Strategy 1: paragraph-level font size
        if para.font.size is not None:
            results["paragraph_font_size_pt"] = emu_to_pt(para.font.size)
            if results["source"] == "none":
                results["source"] = "paragraph.font.size"

        # Strategy 2: run-level font size (from actual text runs)
        for run in para.runs:
            if results["sample_text"] == "" and run.text.strip():
                results["sample_text"] = run.text.strip()[:40]
            if run.font.size is not None:
                results["run_font_size_pt"] = emu_to_pt(run.font.size)
                if results["source"] in ("none", "paragraph.font.size"):
                    results["source"] = "run.font.size"

        # Strategy 3: access defRPr (default run properties) via XML
        # This is the formatting defined on the paragraph even when no runs exist.
        pPr = para._p.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
        )
        if pPr is not None:
            defRPr = pPr.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}defRPr"
            )
            if defRPr is not None and defRPr.get("sz") is not None:
                # sz is in hundredths of a point
                size_hundredths = int(defRPr.get("sz"))
                results["defRPr_size_pt"] = size_hundredths / 100
                if results["source"] == "none":
                    results["source"] = "defRPr (XML)"

    # If nothing found at paragraph/run level, also check endParaRPr
    # (end-of-paragraph run properties, often carries the "intended" font)
    if results["source"] == "none":
        for para in tf.paragraphs:
            endParaRPr = para._p.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}endParaRPr"
            )
            if endParaRPr is not None and endParaRPr.get("sz") is not None:
                size_hundredths = int(endParaRPr.get("sz"))
                results["defRPr_size_pt"] = size_hundredths / 100
                results["source"] = "endParaRPr (XML)"
                break

    return results


def find_pptx_file():
    """Find a .pptx file — check uploads first, then masters."""
    search_dirs = [
        Path(__file__).parent / "slidebuddy" / "data" / "uploads",
        Path(__file__).parent / "slidebuddy" / "data" / "masters",
    ]
    for search_dir in search_dirs:
        if not search_dir.exists():
            continue
        for pptx_path in search_dir.rglob("*.pptx"):
            # Skip temp files (start with ~$)
            if pptx_path.name.startswith("~$"):
                continue
            return str(pptx_path)
    return None


def main():
    # Determine file path
    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
    else:
        pptx_path = find_pptx_file()

    if not pptx_path:
        print("No .pptx file found. Pass a path as argument or place one in uploads/.")
        sys.exit(1)

    print(f"Analyzing: {pptx_path}")
    print("=" * 90)

    prs = Presentation(pptx_path)

    for layout_idx, layout in enumerate(prs.slide_layouts):
        print(f"\n{'─' * 90}")
        print(f"LAYOUT {layout_idx}: \"{layout.name}\"")
        print(f"{'─' * 90}")

        if not layout.placeholders:
            print("  (no placeholders)")
            continue

        for ph in layout.placeholders:
            ph_type = str(ph.placeholder_format.type)
            idx = ph.placeholder_format.idx
            w = emu_to_inches(ph.width)
            h = emu_to_inches(ph.height)

            font_info = get_font_info(ph)

            # Build font size display
            font_parts = []
            if font_info["paragraph_font_size_pt"] is not None:
                font_parts.append(f"para={font_info['paragraph_font_size_pt']}pt")
            if font_info["run_font_size_pt"] is not None:
                font_parts.append(f"run={font_info['run_font_size_pt']}pt")
            if font_info["defRPr_size_pt"] is not None:
                font_parts.append(f"defRPr={font_info['defRPr_size_pt']}pt")

            font_display = " | ".join(font_parts) if font_parts else "None (inherited from theme/master)"

            print(f"  [{idx:>2}] {ph.name:<35} {ph_type}")
            print(f"       Size: {w}\" x {h}\"  |  Paragraphs: {font_info['paragraph_count']}")
            print(f"       Font: {font_display}  (source: {font_info['source']})")
            if font_info["sample_text"]:
                print(f"       Text: \"{font_info['sample_text']}\"")

    # Summary: theme/master-level defaults
    print(f"\n{'=' * 90}")
    print("NOTES ON FONT INHERITANCE:")
    print("  - 'None (inherited)' means the size comes from the slide master or theme")
    print("  - 'defRPr' = default run properties defined on the paragraph in XML")
    print("  - 'endParaRPr' = end-of-paragraph run properties (fallback)")
    print("  - 'paragraph.font.size' = explicit size on paragraph level")
    print("  - 'run.font.size' = explicit size on a text run")
    print("  - python-pptx does NOT resolve inherited theme fonts automatically")
    print("  - To get the effective font, you must walk: run -> paragraph -> layout -> master -> theme")


if __name__ == "__main__":
    main()
