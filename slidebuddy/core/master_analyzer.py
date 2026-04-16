"""Analyze a PowerPoint master file and extract layout structures.

Parses each slide layout to extract placeholder names, types, positions,
and generates a structured description that an LLM can use to suggest
template names, descriptions, and generation prompts.
"""

import json
import re

# Generic OBJECT placeholder names that are decorative, not content
_GENERIC_OBJECT_PATTERNS = re.compile(
    r"^(Inhaltsplatzhalter|Content Placeholder|Platzhalter|Placeholder|Shadow|Rectangle|Oval|"
    r"Freeform|TextBox|Text Box|Group|Picture)\b",
    re.IGNORECASE,
)


def _is_named_object(name: str) -> bool:
    """Check if an OBJECT placeholder has a descriptive (user-assigned) name.

    Returns True for names like 'count1_placeholder', 'conclusion_placeholder'.
    Returns False for generic names like 'Inhaltsplatzhalter 4', 'Shadow'.
    """
    return not bool(_GENERIC_OBJECT_PATTERNS.match(name))


def analyze_master(file_path: str) -> list[dict]:
    """Parse a PPTX master and return structured layout info.

    Returns a list of dicts, one per layout:
    {
        "layout_index": int,
        "layout_name": str,
        "placeholders": [
            {
                "idx": int,
                "name": str,
                "type": str,       # TITLE, BODY, PICTURE, OBJECT
                "is_content": bool, # True for fillable text/image placeholders
                "position": {"left": float, "top": float},  # inches
                "size": {"width": float, "height": float},  # inches
            }
        ],
        "content_placeholders": [...],  # Only fillable ones (TITLE, BODY, PICTURE)
        "structure_summary": str,       # Human-readable summary
    }
    """
    from pptx import Presentation
    prs = Presentation(file_path)
    layouts = []

    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            ph_type = str(ph.placeholder_format.type).split("(")[0].strip()
            # Map enum to readable name
            type_map = {
                "TITLE (1)": "TITLE",
                "BODY (2)": "BODY",
                "PICTURE (18)": "PICTURE",
                "OBJECT (7)": "OBJECT",
            }
            readable_type = type_map.get(str(ph.placeholder_format.type), str(ph.placeholder_format.type))

            # OBJECT placeholders: only include if they have a descriptive name
            # (not generic PowerPoint defaults like "Inhaltsplatzhalter 4")
            if readable_type == "OBJECT":
                is_content = _is_named_object(ph.name)
            else:
                is_content = readable_type in ("TITLE", "BODY", "PICTURE")

            placeholders.append({
                "idx": ph.placeholder_format.idx,
                "name": ph.name,
                "type": readable_type,
                "is_content": is_content,
                "position": {
                    "left": round(ph.left / 914400, 2),
                    "top": round(ph.top / 914400, 2),
                },
                "size": {
                    "width": round(ph.width / 914400, 2),
                    "height": round(ph.height / 914400, 2),
                },
            })

        content_phs = [p for p in placeholders if p["is_content"]]

        # Build human-readable summary
        summary_parts = []
        title_phs = [p for p in content_phs if p["type"] == "TITLE"]
        body_phs = [p for p in content_phs if p["type"] == "BODY"]
        image_phs = [p for p in content_phs if p["type"] == "PICTURE"]
        if title_phs:
            summary_parts.append(f"{len(title_phs)} Titel")
        if body_phs:
            names = [p["name"] for p in body_phs]
            summary_parts.append(f"{len(body_phs)} Textfelder ({', '.join(names)})")
        if image_phs:
            summary_parts.append(f"{len(image_phs)} Bilder")

        layouts.append({
            "layout_index": i,
            "layout_name": layout.name,
            "placeholders": placeholders,
            "content_placeholders": content_phs,
            "structure_summary": " | ".join(summary_parts) if summary_parts else "Leer",
        })

    return layouts


def generate_template_key(layout_name: str) -> str:
    """Generate a valid template key from layout name."""
    key = layout_name.lower().strip()
    # Replace German umlauts
    key = key.replace("ä", "ae").replace("ö", "oe").replace("ü", "ue").replace("ß", "ss")
    # Replace non-alphanumeric with underscore
    key = re.sub(r"[^a-z0-9]+", "_", key)
    key = key.strip("_")
    return key or "unknown"


def estimate_text_capacity(placeholder: dict) -> dict:
    """Estimate text capacity for a placeholder based on box size and name heuristics.

    Since font sizes are typically inherited from the theme and not directly
    readable via python-pptx, we estimate based on placeholder name patterns
    and box dimensions.

    If size data is unavailable (legacy imports), uses conservative defaults
    based on placeholder name patterns.

    Returns dict with: estimated_font_pt, max_lines, max_words, text_hint
    """
    size = placeholder.get("size", {})
    w = size.get("width", 0)
    h = size.get("height", 0)
    name_lower = placeholder.get("name", "").lower()

    # Estimate font size from placeholder name pattern
    if placeholder.get("type") == "TITLE":
        font_pt = 28
    elif "subheading" in name_lower or "heading" in name_lower:
        font_pt = 18
    elif "quote" in name_lower or "statement" in name_lower:
        font_pt = 22
    elif "person" in name_lower or "count" in name_lower:
        font_pt = 14
    elif "bullet" in name_lower:
        font_pt = 14
    elif "conclusion" in name_lower or "fazit" in name_lower or "bridge" in name_lower:
        font_pt = 14
    else:
        # Default body text
        font_pt = 14

    # If we have actual dimensions, calculate from geometry
    if w > 0 and h > 0:
        char_width_inches = (font_pt * 0.6) / 72
        line_height_inches = (font_pt * 1.4) / 72

        chars_per_line = max(1, int(w / char_width_inches)) if char_width_inches > 0 else 20
        max_lines = max(1, int(h / line_height_inches)) if line_height_inches > 0 else 2
        max_words = max(1, int(chars_per_line * max_lines / 5.5))
    else:
        # Fallback: conservative defaults based on field name pattern
        max_lines, max_words = _fallback_capacity(name_lower, placeholder.get("type", ""))

    # Generate human-readable hint
    if max_lines <= 1:
        hint = f"max 1 Zeile, {min(max_words, 8)} Woerter"
    elif max_lines <= 3:
        hint = f"max {max_lines} Zeilen, {max_words} Woerter"
    else:
        hint = f"max {max_lines} Zeilen, {max_words} Woerter"

    return {
        "estimated_font_pt": font_pt,
        "max_lines": max_lines,
        "max_words": max_words,
        "text_hint": hint,
    }


def _fallback_capacity(name_lower: str, ph_type: str) -> tuple[int, int]:
    """Conservative word/line limits when placeholder dimensions are unavailable.

    Returns (max_lines, max_words) based on the placeholder name pattern.
    These are intentionally conservative — it's better to generate slightly
    too little text than to overflow the slide.
    """
    if ph_type == "TITLE":
        return 1, 8
    if "bridge" in name_lower:
        return 1, 3  # Bridge: 1-3 keywords only
    if "conclusion" in name_lower or "fazit" in name_lower:
        return 1, 10  # One short sentence
    if "subheading" in name_lower or "heading" in name_lower:
        return 1, 5  # Short heading
    if "quote" in name_lower or "statement" in name_lower:
        return 2, 20  # A punchy quote
    if "person" in name_lower or "count" in name_lower:
        return 1, 5  # Name or number
    if "bullet" in name_lower:
        return 2, 15  # Short bullet point
    # Default text field: ~2-3 sentences
    return 4, 25


def build_content_schema(content_placeholders: list[dict]) -> dict:
    """Build a JSON content schema from the content placeholders.

    Skips TITLE (goes to top-level title) and PICTURE (user inserts manually).
    Only includes text-fillable placeholders.
    """
    schema = {}
    for ph in content_placeholders:
        field_name = ph["name"]
        if ph["type"] == "TITLE":
            continue
        elif ph["type"] == "PICTURE":
            # Images are inserted manually by the user — skip
            continue
        else:
            capacity = estimate_text_capacity(ph)
            schema[field_name] = f"str — {capacity['text_hint']}"
    return schema


def build_generation_prompt(
    layout_name: str,
    content_placeholders: list[dict],
    content_schema: dict,
) -> str:
    """Build a generation prompt that matches the prompt_assembler format.

    Produces HINWEISE text that tells the LLM exactly how to fill the content
    JSON for this layout, referencing the actual placeholder field names
    and their text capacity constraints.
    """
    # Classify placeholders
    body_phs = [p for p in content_placeholders if p["type"] == "BODY"]
    image_phs = [p for p in content_placeholders if p["type"] == "PICTURE"]

    # Detect structural patterns
    subheading_phs = [p for p in body_phs if "subheading" in p["name"].lower() or "heading" in p["name"].lower()]
    text_phs = [p for p in body_phs if "text" in p["name"].lower() and "subheading" not in p["name"].lower() and "heading" not in p["name"].lower()]
    bullet_phs = [p for p in body_phs if "bullet" in p["name"].lower()]
    quote_phs = [p for p in body_phs if "quote" in p["name"].lower() or "statement" in p["name"].lower()]
    conclusion_phs = [p for p in body_phs if "conclusion" in p["name"].lower() or "fazit" in p["name"].lower()]
    person_phs = [p for p in body_phs if "person" in p["name"].lower()]
    bridge_phs = [p for p in body_phs if "bridge" in p["name"].lower()]

    lines = []

    # Count repeating groups (e.g. subheading1+text1, subheading2+text2)
    num_groups = max(len(subheading_phs), len(text_phs), 1)
    has_groups = len(subheading_phs) > 1 or len(text_phs) > 1

    # General structure hint
    if has_groups:
        lines.append(f"- Dieses Layout hat {num_groups} inhaltliche Abschnitte die parallel befuellt werden muessen")

    # Field-specific hints with text capacity
    if subheading_phs and text_phs:
        lines.append("- Alle Abschnitte sollten ungefaehr gleich lang sein")
        lines.append("- Ueberschriften sollten parallel formuliert sein (gleicher Satzbau)")

    if bullet_phs:
        lines.append(f"- {len(bullet_phs)} Stichpunkt-Felder — jedes mit einem praegnanten Satz fuellen")

    if quote_phs:
        lines.append("- Zitatfeld mit einem aussagekraeftigen Zitat oder These fuellen")

    if person_phs:
        lines.append("- Personenfeld mit Name und ggf. Rolle/Titel fuellen")

    if conclusion_phs:
        lines.append("- Fazit/Conclusion-Feld: eine kurze Kernaussage, maximal 1 Satz")

    if bridge_phs:
        lines.append("- Bridge-Feld: NUR 1-3 Schlagwoerter (z.B. 'Kernprozesse', 'Ausblick'), KEIN ganzer Satz")

    # Per-field text capacity hints
    lines.append("")
    lines.append("TEXTLAENGEN pro Feld (Platz auf der Folie):")
    for ph in content_placeholders:
        if ph["type"] in ("TITLE", "PICTURE"):
            continue
        capacity = estimate_text_capacity(ph)
        size = ph.get("size", {})
        if size.get("width") and size.get("height"):
            dims = f' ({size["width"]:.1f}x{size["height"]:.1f} Zoll, ~{capacity["estimated_font_pt"]}pt)'
        else:
            dims = ""
        lines.append(f'  "{ph["name"]}": {capacity["text_hint"]}{dims}')

    # JSON output reminder — only text fields (no images)
    schema_fields = list(content_schema.keys())
    if schema_fields:
        lines.append("")
        field_list = ", ".join(f'"{f}"' for f in schema_fields[:8])
        if len(schema_fields) > 8:
            field_list += f" ... ({len(schema_fields)} Felder gesamt)"
        lines.append(f"- Dein JSON-content muss exakt diese Felder enthalten: {field_list}")
        lines.append("- Verwende NUR die genannten Feldnamen, keine eigenen erfinden")
        if image_phs:
            lines.append(f"- {len(image_phs)} Bildplatzhalter vorhanden — diese werden vom Nutzer manuell befuellt, NICHT im JSON ausgeben")

    # Fallback if nothing was detected
    if not lines:
        lines.append("- Befuelle alle Felder im JSON-Schema mit passendem Inhalt")
        lines.append("- Halte Texte praegnant und passend zum Folienthema")

    return "\n".join(lines)


def build_llm_analysis_prompt(layouts: list[dict]) -> str:
    """Build a prompt for LLM to analyze layouts and suggest template metadata.

    The LLM generates purpose, display_name, description, and generation_prompt
    based on the layout structure — replacing the old hardcoded purpose_map.
    """
    layout_descriptions = []
    for layout in layouts:
        if not layout["content_placeholders"]:
            continue  # Skip empty layouts

        phs = []
        for ph in layout["content_placeholders"]:
            capacity = estimate_text_capacity(ph)
            size = ph.get("size", {})
            if size.get("width") and size.get("height"):
                dims = f"{size['width']:.1f}x{size['height']:.1f} Zoll, "
            else:
                dims = ""
            phs.append(
                f"  - {ph['name']} ({ph['type']}, "
                f"{dims}"
                f"~{capacity['estimated_font_pt']}pt, {capacity['text_hint']})"
            )

        layout_descriptions.append(
            f"Layout {layout['layout_index']}: \"{layout['layout_name']}\"\n"
            f"Placeholders:\n" + "\n".join(phs)
        )

    layouts_text = "\n\n".join(layout_descriptions)

    return f"""Analysiere die folgenden PowerPoint-Folienlayouts und erstelle fuer jedes Layout:

1. **display_name**: Ein kurzer, klarer deutscher Name (z.B. "Vergleich Links/Rechts")
2. **description**: 1-2 Saetze Beschreibung wofuer das Layout ideal ist
3. **purpose**: Ein Satz der das ZIEL dieser Folie beschreibt (z.B. "zwei Aspekte gegenueberstellen und Unterschiede herausarbeiten"). Leite dies aus Layout-Name und Placeholder-Struktur ab.
4. **generation_prompt**: Ein Prompt-Absatz (3-5 Saetze) der dem LLM erklaert wie Content fuer dieses Layout generiert werden soll. Beschreibe welche Felder befuellt werden muessen und wie der Inhalt strukturiert sein soll. Beachte die Textlaengen-Limits der Felder.

LAYOUTS:

{layouts_text}

Antworte als JSON-Array. Jedes Element hat:
{{
    "layout_index": int,
    "display_name": "str",
    "description": "str",
    "purpose": "str",
    "generation_prompt": "str"
}}

Ueberspringe leere Layouts (ohne Content-Placeholders).
Antworte NUR mit dem JSON-Array, ohne Markdown-Fences."""


# ---------------------------------------------------------------------------
# Orchestration: re-analyze an existing master and update DB schemas
# ---------------------------------------------------------------------------

import logging
import sqlite3

_logger = logging.getLogger(__name__)


def reanalyze_master_templates(conn: sqlite3.Connection, master_id: str) -> int:
    """Re-analyze an existing master's PPTX to update content schemas with word limits.

    Reads the PPTX file, extracts placeholder dimensions, recalculates
    content_schema (with max_words hints) and generation_prompt (with
    TEXTLAENGEN), then updates each template in the DB.

    Returns the number of templates updated.
    """
    from slidebuddy.db.queries import get_slide_master, get_templates_for_master, update_master_template

    master = get_slide_master(conn, master_id)
    if not master:
        raise ValueError(f"Master {master_id} not found")

    # Re-analyze the PPTX file
    layouts = analyze_master(master.file_path)
    layout_by_index = {l["layout_index"]: l for l in layouts}

    templates = get_templates_for_master(conn, master_id)
    updated_count = 0

    for tpl in templates:
        layout = layout_by_index.get(tpl.layout_index)
        if not layout or not layout["content_placeholders"]:
            continue

        content_phs = layout["content_placeholders"]

        # Rebuild content_schema with word limit hints
        new_schema = build_content_schema(content_phs)

        # Rebuild generation_prompt with TEXTLAENGEN section
        new_gen_prompt = build_generation_prompt(
            layout["layout_name"], content_phs, new_schema,
        )

        # Preserve the existing ZIEL line from the old generation_prompt
        if tpl.generation_prompt:
            import re as _re
            ziel_match = _re.search(r'(- ZIEL[^\n]+)', tpl.generation_prompt)
            if ziel_match:
                new_gen_prompt = ziel_match.group(1) + "\n" + new_gen_prompt

        # Update placeholder_schema to include size data
        new_ph_schema = json.dumps(
            [
                {
                    "name": ph["name"],
                    "type": ph["type"],
                    "idx": ph["idx"],
                    "size": ph.get("size", {}),
                }
                for ph in content_phs
            ]
        )

        tpl.placeholder_schema = new_ph_schema
        tpl.content_schema = json.dumps(new_schema)
        tpl.generation_prompt = new_gen_prompt
        update_master_template(conn, tpl)
        updated_count += 1

        _logger.info(
            "Updated template '%s': %d fields with word limits",
            tpl.template_key, len(new_schema),
        )

    # Clear prompt caches so new schemas are used immediately
    from slidebuddy.llm.prompt_assembler import clear_template_cache
    clear_template_cache()

    # Clear word limit cache in slide generation
    from slidebuddy.core.slide_generation import _word_limit_cache
    _word_limit_cache.clear()

    return updated_count
