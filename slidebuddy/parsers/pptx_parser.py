from pathlib import Path

def parse_pptx(file_path: Path) -> str:
    """Extract text and speaker notes from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, 1):
        parts = [f"--- Slide {slide_num} ---"]

        # Extract text from all shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"\nSPRECHERNOTIZEN:\n{notes}")

        if len(parts) > 1:  # More than just the header
            slides_text.append("\n".join(parts))

    return "\n\n".join(slides_text)


def parse_pptx_as_slides(file_path: Path) -> list[dict]:
    """Parse PPTX into individual slide dicts for RAG ingestion."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    slides = []

    for slide_num, slide in enumerate(prs.slides, 1):
        title = ""
        content_parts = []
        notes = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if shape.shape_id == slide.shapes.title.shape_id if slide.shapes.title else False:
                    title = text
                else:
                    content_parts.append(text)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        if title or content_parts:
            slides.append({
                "slide_number": slide_num,
                "title": title,
                "content": "\n".join(content_parts),
                "speaker_notes": notes,
                "source_file": file_path.name,
            })

    return slides
