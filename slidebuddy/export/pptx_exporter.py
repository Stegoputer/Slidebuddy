"""Export generated slides to a PowerPoint file using the active master template.

If a master is active, slides are placed on matching layouts with content
filled into the correct placeholders. Without a master, a basic default
presentation is created.
"""

import json
import logging
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)

from slidebuddy.config.defaults import DB_PATH
from slidebuddy.db.migrations import get_connection
from slidebuddy.db.queries import get_active_slide_master, get_active_master_templates


def export_pptx(
    project_name: str,
    gen_slides: dict[int, list[dict]],
    chapters: list | None = None,
) -> bytes:
    """Export generated slides to PPTX bytes.

    Uses the active master template if available, otherwise creates
    a basic presentation.
    """
    conn = get_connection(DB_PATH)
    master = get_active_slide_master(conn)
    master_templates = []
    if master:
        master_templates = get_active_master_templates(conn)
    conn.close()

    if master and master_templates:
        return _export_with_master(project_name, gen_slides, chapters, master, master_templates)
    return _export_default(project_name, gen_slides, chapters)


def _export_with_master(
    project_name: str,
    gen_slides: dict[int, list[dict]],
    chapters: list | None,
    master,
    master_templates: list,
) -> bytes:
    """Export using the master PPTX as base, filling placeholders."""
    prs = Presentation(master.file_path)

    # Remove any existing slides from the master (keep only layouts)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get("r:id")  # noqa: SLF001
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])  # noqa: SLF001

    tpl_by_key = {t.template_key: t for t in master_templates}

    for ch_idx in sorted(gen_slides.keys()):
        for slide_data in gen_slides[ch_idx]:
            template_key = slide_data.get("template_type", "")
            tpl = tpl_by_key.get(template_key)

            if tpl and tpl.layout_index < len(prs.slide_layouts):
                layout = prs.slide_layouts[tpl.layout_index]
                slide = prs.slides.add_slide(layout)
                _fill_placeholders(slide, slide_data, tpl)
            else:
                # Fallback: use first layout
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                _fill_basic(slide, slide_data)

    # Add speaker notes
    for ch_idx in sorted(gen_slides.keys()):
        slide_offset = sum(len(gen_slides.get(i, [])) for i in range(ch_idx))
        for i, slide_data in enumerate(gen_slides[ch_idx]):
            slide_idx = slide_offset + i
            if slide_idx < len(prs.slides):
                notes = slide_data.get("speaker_notes", "")
                if notes:
                    slide = prs.slides[slide_idx]
                    if not slide.has_notes_slide:
                        slide.notes_slide  # noqa: B018 — creates notes slide
                    slide.notes_slide.notes_text_frame.text = notes

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _fill_placeholders(slide, slide_data: dict, tpl):
    """Fill slide placeholders from generated content using the template schema.

    Uses the placeholder_schema from the MasterTemplate to build a direct
    mapping from content field names to placeholder indices, avoiding
    unreliable fuzzy name matching.
    """
    content = slide_data.get("content", {})
    if isinstance(content, str):
        try:
            content = json.loads(content)
        except json.JSONDecodeError:
            content = {}

    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")

    # Build name→idx map from the stored placeholder schema
    # This maps e.g. "Subheading 1 Placeholder" → 10, "Text Placeholder 2" → 11
    name_to_idx: dict[str, int] = {}
    if tpl.placeholder_schema:
        try:
            phs = json.loads(tpl.placeholder_schema)
            for ph in phs:
                name_to_idx[ph["name"]] = ph["idx"]
        except (json.JSONDecodeError, KeyError):
            pass

    # Build idx→value map from content fields using the schema mapping
    idx_to_value: dict[int, str] = {}
    unmatched_fields = []
    for field_name, value in content.items():
        if isinstance(value, str) and field_name in name_to_idx:
            idx_to_value[name_to_idx[field_name]] = value
        elif isinstance(value, str):
            unmatched_fields.append(field_name)

    if unmatched_fields:
        logger.warning(
            "PPTX export: content fields not in placeholder_schema: %s. "
            "Schema names: %s",
            unmatched_fields, list(name_to_idx.keys()),
        )

    logger.debug(
        "PPTX fill: template=%s, schema_map=%s, idx_to_value_keys=%s, content_keys=%s",
        tpl.template_key, name_to_idx, list(idx_to_value.keys()), list(content.keys()),
    )

    # Fill each placeholder on the slide
    for ph in slide.placeholders:
        ph_idx = ph.placeholder_format.idx
        ph_name = ph.name.lower()

        # Title placeholder
        if ph_idx == 0 or "title" in ph_name:
            ph.text = title
            continue

        # Direct index match from schema
        if ph_idx in idx_to_value:
            ph.text = idx_to_value[ph_idx]
            continue

        # Fallback: try fuzzy matching for content fields not in schema
        matched = False
        for field_name, value in content.items():
            if isinstance(value, str) and _names_match(field_name, ph.name):
                ph.text = value
                matched = True
                break

        # Subtitle fallback
        if not matched and ("subtitle" in ph_name or "subheading" in ph_name):
            if subtitle:
                ph.text = subtitle


def _names_match(field_name: str, placeholder_name: str) -> bool:
    """Fuzzy-check if a content field name matches a placeholder name.

    Used as fallback when schema-based matching doesn't find a match.
    """
    fn = field_name.lower().replace("_", " ").replace("-", " ").split()
    pn = placeholder_name.lower().replace("_", " ").replace("-", " ").split()
    # Check if all significant words from the field name appear in the placeholder name
    # Skip common filler words
    skip = {"placeholder", "text", "field"}
    fn_words = [w for w in fn if w not in skip and len(w) > 1]
    pn_words = [w for w in pn if w not in skip and len(w) > 1]
    if not fn_words or not pn_words:
        return False
    return all(any(fw in pw or pw in fw for pw in pn_words) for fw in fn_words)


def _fill_basic(slide, slide_data: dict):
    """Fill a slide with basic title/subtitle when no template match."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = slide_data.get("title", "")
        elif ph.placeholder_format.idx == 1:
            ph.text = slide_data.get("subtitle", "")


# ---------------------------------------------------------------------------
# Default export (no master)
# ---------------------------------------------------------------------------

def _export_default(
    project_name: str,
    gen_slides: dict[int, list[dict]],
    chapters: list | None,
) -> bytes:
    """Create a basic PPTX without a master template."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for ch_idx in sorted(gen_slides.keys()):
        ch_title = ""
        if chapters and ch_idx < len(chapters):
            ch = chapters[ch_idx]
            ch_title = ch.title if hasattr(ch, "title") else ch.get("title", "")

        # Chapter divider slide
        layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(layout)
        slide.placeholders[0].text = ch_title or f"Kapitel {ch_idx + 1}"

        # Content slides
        for slide_data in gen_slides[ch_idx]:
            layout = prs.slide_layouts[1]  # Title + Content
            slide = prs.slides.add_slide(layout)
            slide.placeholders[0].text = slide_data.get("title", "")
            if len(slide.placeholders) > 1:
                content = slide_data.get("content", {})
                if isinstance(content, str):
                    try:
                        content = json.loads(content)
                    except json.JSONDecodeError:
                        pass
                slide.placeholders[1].text = _flatten_content(content)

            # Speaker notes
            notes = slide_data.get("speaker_notes", "")
            if notes:
                if not slide.has_notes_slide:
                    slide.notes_slide  # noqa: B018
                slide.notes_slide.notes_text_frame.text = notes

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _flatten_content(content) -> str:
    """Flatten content dict/str to plain text for default export."""
    if isinstance(content, str):
        return content
    if not isinstance(content, dict):
        return str(content)

    parts = []
    for key, value in content.items():
        if isinstance(value, str) and value.strip():
            parts.append(value)
        elif isinstance(value, list):
            for item in value:
                if isinstance(item, dict):
                    heading = item.get("heading", "")
                    text = item.get("text", "")
                    if heading:
                        parts.append(f"{heading}: {text}")
                    elif text:
                        parts.append(text)
                elif isinstance(item, str):
                    parts.append(item)
        elif isinstance(value, dict):
            heading = value.get("heading", "")
            text = value.get("text", "")
            if heading or text:
                parts.append(f"{heading}: {text}" if heading else text)
    return "\n".join(parts)
